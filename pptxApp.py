import shutil
from sqlalchemy.orm import Session
from pptx import Presentation
import openai
import json
import os
import re
import asyncio
from handle_db import Upload, engine

openai.api_key = os.environ.get('API_KEY')
CONTENT = [
    {"role": "system", "content": "Can you explain the slides in basic english, and provide examples if needed!"}
]
ERROR_MESSAGE = "Something is wrong:"
ENGINE_MODEL = "gpt-3.5-turbo"
WRITE_TO_FILE_MODE = 'w'
UPLOADS_FOLDER = 'uploads'
OUTPUTS_FOLDER = 'outputs'
PROCESSED_FOLDER = 'processed'
EXPLANATION_SAVED = "Explanations saved to"
EXPLANATION_SAVED_ERROR = "Error saving explanations:"
PATH_NOT_FOUND = "the path you provided does not exist."
PROCESS_SLIDE_ERROR = "Error processing slide:"
MOVED_FILE = "Moved file:"
PROCESSING_FILE = "Processing file:"
DONE_STATUS = 'done'
PROCESSING_STATUS = 'processing'
PENDING_STATUS = 'pending'
PROCESS_FILE_ERROR = "Error processing file"
EXPLAINER_STARTED_MESSAGE = "Explainer started."
CHOICES = "choices"
FIRST_ELEMENT = 0


async def parse_presentation(presentation_path):
    """
    This method receives a path for a pptx presentation, checks if the path is found in the operating system, if so,
    parses the data to slides. The method, returns a list of explanations.
    :param: presentation_path: path of a power-point presentation. (String)
    :return: list of explanations. (List of strings)
    """
    # check if path is available
    if not os.path.isfile(presentation_path):
        print(f"{ERROR_MESSAGE} {PATH_NOT_FOUND}")
        return []

    prs = Presentation(presentation_path)
    explanations = []
    for slide in prs.slides:
        explanation = await parse_slide_of_pptx(slide)
        explanations.append(explanation)
    return explanations


async def parse_slide_of_pptx(slide):
    """
    This method receives a single slide, the method parses the data on that slide, it calls another method to get the
    response and return it to the parse_presentation method. It throws an error if an exception occurred.
    :param: slide: A single slide from the power-point. (slide object)
    :return: The explanation if the processing went well, an error, otherwise. (List of Strings)
    """
    try:
        response = await request_completion(slide)
        return response
    except Exception as error:
        error_message = f"{ERROR_MESSAGE} {PROCESS_SLIDE_ERROR} {str(error)}"
        return error_message


def parse_text_of_slide(slide):
    """
    This method receives a slide, and extracts all the extractable text found in that slide. In addition, it cleans
    the data using the strip function.
    :param: slide: A single slide from the power-point. (Slide Object)
    :return: Explanation response. (String)
    """
    slide_text = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                slide_text.append(run.text.strip())
    return " ".join(slide_text)


async def request_completion(slide):
    """
    This method receives a slide object, sends a request to the openai API asking the server to explain the content
    of that slide, the method returns the response from the API.
    :param: slide: Slide: a single slide from the power-point. (Slide Object)
    :return: Response of the API.
    """
    CONTENT.append({"role": "user", "content": parse_text_of_slide(slide)})
    response = await openai.ChatCompletion.acreate(
        model=ENGINE_MODEL,
        messages=CONTENT
    )
    content = response[CHOICES][FIRST_ELEMENT].message.content
    cleaned_content = clean_text(content)
    return cleaned_content


def clean_text(text):
    """
    This method receives the response retrieved from the openai API and cleans it, in other words, gets rid of
    unwanted characters.
    :param: text: Response from API. (string)
    :return: clean Version of the response. (string)
    """
    cleaned_text = re.sub(r"\n", "", text)
    cleaned_text = cleaned_text.encode("ascii", "ignore").decode("utf-8")
    return cleaned_text.strip()


def save_explanations(explanations, file_path):
    """
    This method receives a list of strings, each string representing an explanation. It creates a JSON file and appends
    the explanations to the file, using the original file name.
    :param: explanations: List of explanations retrieved from the API. (List of strings)
    :param: file_path: Path of the original file. (String)
    :return:
    """
    file_name = os.path.basename(file_path)
    presentation_name, extension = os.path.splitext(file_name)
    output_file = os.path.join(OUTPUTS_FOLDER, f"{presentation_name}.json")
    slide_explanations = {}

    for slide_num, explanation in enumerate(explanations, start=1):
        slide_key = f"slide{slide_num}"
        slide_explanations[slide_key] = explanation

    try:
        if not os.path.exists(OUTPUTS_FOLDER):
            os.makedirs(OUTPUTS_FOLDER)

        with open(output_file, WRITE_TO_FILE_MODE) as file:
            json.dump(slide_explanations, file, indent=4)
        print(f"{EXPLANATION_SAVED} {output_file}")
    except IOError as error:
        print(f"{EXPLANATION_SAVED_ERROR} {str(error)}")


def move_file(file_path, destination_folder):
    """
    This method receives a file path, and a destination folder. The method moves the processed folder from the
    'uploads' folder to the 'processed' folder in order to prevent the explainer to scan old files that have been
    already processed
    :param file_path: path of a file (string)
    :param destination_folder: name of the processed folder (string)
    :return:
    """
    file_name = os.path.basename(file_path)
    destination_path = os.path.join(destination_folder, file_name)

    if not os.path.exists(destination_folder):
        os.makedirs(destination_folder)

    shutil.move(file_path, destination_path)
    print(f"{MOVED_FILE} {file_path} to {destination_path}")


async def process_file(file_path, file_processing):
    """
    This method receives a file_path, calls the needed functions to explain the contents of the power-point, save and
    moves the processed file. it throws an exception if it could not process a file.
    :param: file_processing: upload object representing the file being processed (Upload)
    :param: file_path: a file path from the 'uploads' folder (string)
    :return:
    """
    print(f"{PROCESSING_FILE} {file_path}")
    try:
        explanations = await parse_presentation(file_path)
        save_explanations(explanations, file_path)
        move_file(file_path, PROCESSED_FOLDER)

        file_processing.set_file_status(DONE_STATUS)
        file_processing.set_upload_finish_time()
    except Exception as error:
        error_message = f"{ERROR_MESSAGE} {PROCESS_FILE_ERROR} {str(error)}"
        print(error_message)


async def main_loop():
    """
    This method keeps running in an infinite loop, each iteration it searches for all pending files using sql queries
    for each file found, it retrieves it from the 'uploads' folder, updates its status to processing and then
    waits for it to be processed.

    :return:
    """
    while True:
        await asyncio.sleep(10)
        with Session(engine) as session:
            pending_files = session.query(Upload).filter_by(status=PENDING_STATUS).all()
            if pending_files:
                for pending_file in pending_files:
                    pending_file.set_file_status(PROCESSING_STATUS)
                    session.add(pending_file)
                    session.commit()
                    file_path = pending_file.get_upload_path()
                    await process_file(file_path, pending_file)
                    session.add(pending_file)
                    session.commit()


if __name__ == "__main__":
    print(EXPLAINER_STARTED_MESSAGE)
    asyncio.run(main_loop())
